"""
Generate AgentOps Presentation - Leveraging Microsoft Foundry & Agent 365
Run: pip install python-pptx && python scripts/generate_agentops_pptx.py
Output: docs/presentations/AgentOps-Microsoft-Foundry-Agent365.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

# -- Color Palette (Microsoft-inspired) --
DARK_BLUE = RGBColor(0x00, 0x78, 0xD4)       # Microsoft Blue
LIGHT_BLUE = RGBColor(0x50, 0xE6, 0xFF)       # Accent Cyan
DARK_BG = RGBColor(0x1B, 0x1B, 0x1B)          # Near-Black
CHARCOAL = RGBColor(0x2D, 0x2D, 0x2D)         # Dark Gray
MID_GRAY = RGBColor(0x60, 0x60, 0x60)         # Mid Gray
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)       # Light Gray
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0xB2, 0x94)             # Teal Green
PURPLE = RGBColor(0x88, 0x61, 0xC4)            # Purple
ORANGE = RGBColor(0xFF, 0x8C, 0x00)            # Dark Orange
RED_ACCENT = RGBColor(0xE7, 0x4C, 0x3C)        # Red
YELLOW = RGBColor(0xFF, 0xB9, 0x00)            # Gold


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, shape_type, left, top, width, height, fill_color, line_color=None):
    """Add a shape with fill and optional line."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_to_shape(shape, text, font_size=12, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER):
    """Add formatted text to a shape."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False,
                color=WHITE, alignment=PP_ALIGN.LEFT):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_bullet_text(tf, text, font_size=14, color=WHITE, bold=False, level=0):
    """Add a bullet paragraph to a text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.level = level
    p.space_before = Pt(4)
    p.space_after = Pt(2)
    return p


def add_connector_line(slide, x1, y1, x2, y2, color=LIGHT_BLUE, width=2):
    """Add a simple line connector."""
    connector = slide.shapes.add_connector(
        1, x1, y1, x2, y2  # MSO_CONNECTOR.STRAIGHT = 1
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def add_arrow_shape(slide, left, top, width, height, color=DARK_BLUE):
    """Add a right arrow shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def slide_title_bar(slide, title_text, subtitle_text=None):
    """Add a consistent title bar across content slides."""
    # Title background strip
    bar = add_shape(slide, MSO_SHAPE.RECTANGLE,
                    Inches(0), Inches(0), Inches(13.33), Inches(1.2), DARK_BLUE)
    add_text_to_shape(bar, "", font_size=1)

    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(11), Inches(0.6),
                title_text, font_size=28, bold=True, color=WHITE)

    if subtitle_text:
        add_textbox(slide, Inches(0.6), Inches(0.7), Inches(11), Inches(0.4),
                    subtitle_text, font_size=14, bold=False, color=LIGHT_BLUE)

    # Accent line
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(1.2), Inches(13.33), Inches(0.04), LIGHT_BLUE)


def slide_footer(slide, page_num, total):
    """Add footer with page number."""
    add_textbox(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
                f"{page_num} / {total}", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


# ========================== BUILD PRESENTATION ==========================

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

TOTAL_SLIDES = 15

# ============ SLIDE 1: Title Slide ============
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

# Large hero gradient block
hero = add_shape(slide, MSO_SHAPE.RECTANGLE,
                 Inches(0), Inches(0), Inches(13.33), Inches(4.5), DARK_BLUE)
add_text_to_shape(hero, "", font_size=1)

# Decorative element - side accent
add_shape(slide, MSO_SHAPE.RECTANGLE,
          Inches(0), Inches(0), Inches(0.15), Inches(7.5), LIGHT_BLUE)

# Title
add_textbox(slide, Inches(1.0), Inches(1.0), Inches(11), Inches(1.2),
            "AgentOps", font_size=52, bold=True, color=WHITE)

add_textbox(slide, Inches(1.0), Inches(2.2), Inches(11), Inches(1.0),
            "Operationalizing AI Agents with Microsoft Foundry & Agent 365",
            font_size=24, bold=False, color=LIGHT_BLUE)

# Decorative circles
for i, (x, clr) in enumerate([(Inches(10.5), PURPLE), (Inches(11.3), GREEN), (Inches(12.1), ORANGE)]):
    add_shape(slide, MSO_SHAPE.OVAL, x, Inches(3.6), Inches(0.5), Inches(0.5), clr)

# Subtitle area
add_textbox(slide, Inches(1.0), Inches(5.2), Inches(8), Inches(0.5),
            "Building, Deploying & Governing Enterprise AI Agents at Scale",
            font_size=18, bold=True, color=WHITE)

add_textbox(slide, Inches(1.0), Inches(6.0), Inches(8), Inches(0.5),
            "March 2026  |  Practitioner-Level Presentation",
            font_size=14, bold=False, color=MID_GRAY)

slide_footer(slide, 1, TOTAL_SLIDES)


# ============ SLIDE 2: Agenda ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
slide_title_bar(slide, "Agenda", "What We Will Cover Today")

agenda_items = [
    ("01", "What is AgentOps?", "The discipline of operationalizing AI agents end-to-end"),
    ("02", "Microsoft Foundry", "The agent factory - build, deploy, observe, govern"),
    ("03", "Agent 365", "Enterprise agent lifecycle management across M365"),
    ("04", "Agent Architecture", "Components: Models, Tools, Orchestration, Trust"),
    ("05", "AgentOps Lifecycle", "Design -> Build -> Deploy -> Monitor -> Optimize"),
    ("06", "Foundry Agent Service", "Deep dive into agent runtime capabilities"),
    ("07", "Workflows & Orchestration", "Multi-agent coordination patterns"),
    ("08", "Agent 365 Governance", "Identity, registry, access control, security"),
    ("09", "Observability & Evaluation", "Tracing, metrics, evaluations, feedback loops"),
    ("10", "Enterprise Integration", "Connecting agents to business systems"),
    ("11", "Security & Trust", "RBAC, content safety, data protection"),
    ("12", "Getting Started", "Quickstart roadmap and next steps"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    row = i // 3
    col = i % 3
    x = Inches(0.6 + col * 4.1)
    y = Inches(1.6 + row * 1.35)

    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, y, Inches(3.8), Inches(1.15), CHARCOAL, DARK_BLUE)

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{num}  {title}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.LEFT

    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.LEFT

slide_footer(slide, 2, TOTAL_SLIDES)


# ============ SLIDE 3: What is AgentOps? ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
slide_title_bar(slide, "What is AgentOps?", "The Discipline of Operationalizing AI Agents")

# Definition block
defn = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                 Inches(0.6), Inches(1.6), Inches(12), Inches(1.2), CHARCOAL, DARK_BLUE)
tf = defn.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("AgentOps is the set of practices, tools, and governance frameworks "
          "for managing the full lifecycle of AI agents -- from design and development "
          "through deployment, monitoring, evaluation, and continuous improvement in production.")
p.font.size = Pt(14)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

# Three pillars
pillars = [
    ("Build", "Design, develop, test\nagents with models,\ntools & instructions", DARK_BLUE, "B"),
    ("Deploy", "Ship agents to\nproduction with CI/CD,\nsecurity & compliance", GREEN, "D"),
    ("Operate", "Monitor, evaluate,\niterate agents with\nobservability & feedback", PURPLE, "O"),
]

for i, (title, desc, clr, letter) in enumerate(pillars):
    x = Inches(1.2 + i * 3.8)
    y = Inches(3.2)

    # Circle icon
    circle = add_shape(slide, MSO_SHAPE.OVAL, x + Inches(1.1), y, Inches(1.0), Inches(1.0), clr)
    add_text_to_shape(circle, letter, font_size=28, bold=True, color=WHITE)

    # Label
    add_textbox(slide, x, y + Inches(1.2), Inches(3.3), Inches(0.5),
                title, font_size=20, bold=True, color=clr, alignment=PP_ALIGN.CENTER)

    # Description
    add_textbox(slide, x, y + Inches(1.7), Inches(3.3), Inches(1.2),
                desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Arrows between pillars
for i in range(2):
    x = Inches(4.3 + i * 3.8)
    add_arrow_shape(slide, x, Inches(3.3), Inches(0.8), Inches(0.4), LIGHT_BLUE)

# Bottom callout
callout_text = ("Think of it as DevOps for AI Agents -- applying operational rigor to intelligent, "
                "autonomous systems that reason, decide, and act.")
add_textbox(slide, Inches(0.6), Inches(6.3), Inches(12), Inches(0.6),
            callout_text, font_size=13, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

slide_footer(slide, 3, TOTAL_SLIDES)


# ============ SLIDE 4: Why AgentOps Matters ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
slide_title_bar(slide, "Why AgentOps Matters", "Challenges of Moving AI Agents to Production")

# Problem Statement
problems = [
    ("LLMs Drift & Hallucinate", "Without guardrails, agent outputs degrade over time and produce unreliable results"),
    ("No Visibility", "Teams cannot see what agents decide, which tools they call, or why they fail"),
    ("Security Gaps", "Agents without identity management can be exploited or access unauthorized data"),
    ("No Governance", "Shadow agents proliferate -- untracked, unmanaged, non-compliant"),
    ("Scale Challenges", "Manual orchestration breaks down with multiple agents across business units"),
    ("Trust Deficit", "Stakeholders cannot trust agents they cannot observe, audit, or control"),
]

for i, (title, desc) in enumerate(problems):
    row = i // 2
    col = i % 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(1.6 + row * 1.7)

    # Icon placeholder
    icon_colors = [RED_ACCENT, ORANGE, PURPLE, YELLOW, DARK_BLUE, GREEN]
    icon = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                     x, y, Inches(0.6), Inches(0.6), icon_colors[i])
    add_text_to_shape(icon, str(i + 1), font_size=16, bold=True, color=WHITE)

    add_textbox(slide, x + Inches(0.8), y, Inches(5), Inches(0.4),
                title, font_size=16, bold=True, color=WHITE)
    add_textbox(slide, x + Inches(0.8), y + Inches(0.4), Inches(5), Inches(0.6),
                desc, font_size=11, color=LIGHT_GRAY)

# Bottom banner
banner = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                   Inches(0.6), Inches(6.4), Inches(12), Inches(0.7), DARK_BLUE)
add_text_to_shape(banner,
                  "AgentOps addresses these challenges with structure, tooling & governance",
                  font_size=14, bold=True, color=WHITE)

slide_footer(slide, 4, TOTAL_SLIDES)


# ============ SLIDE 5: Microsoft Foundry Overview ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
slide_title_bar(slide, "Microsoft Foundry", "The Agent Factory -- Build, Deploy, Observe, Govern")

# Central hub diagram
# Center circle
center = add_shape(slide, MSO_SHAPE.OVAL,
                   Inches(5.4), Inches(3.0), Inches(2.5), Inches(2.5), DARK_BLUE, LIGHT_BLUE)
add_text_to_shape(center, "Foundry\nAgent\nService", font_size=16, bold=True, color=WHITE)

# Surrounding components
components = [
    ("AI Models", Inches(2.0), Inches(1.8), GREEN),
    ("Tools &\nFrameworks", Inches(9.2), Inches(1.8), PURPLE),
    ("Governance\n& Compliance", Inches(9.2), Inches(5.0), ORANGE),
    ("Orchestration", Inches(2.0), Inches(5.0), DARK_BLUE),
]

for label, x, y, clr in components:
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, y, Inches(2.0), Inches(1.2), clr)
    add_text_to_shape(box, label, font_size=12, bold=True, color=WHITE)

# Connection lines
connections = [
    (Inches(4.0), Inches(2.4), Inches(5.4), Inches(3.5)),
    (Inches(9.2), Inches(2.4), Inches(7.9), Inches(3.5)),
    (Inches(9.2), Inches(5.6), Inches(7.9), Inches(5.0)),
    (Inches(4.0), Inches(5.6), Inches(5.4), Inches(5.0)),
]
for x1, y1, x2, y2 in connections:
    add_connector_line(slide, x1, y1, x2, y2, LIGHT_BLUE, 2)

# Bottom description
add_textbox(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
            "Foundry combines models, tools, frameworks & governance into a unified system for building intelligent agents",
            font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

slide_footer(slide, 5, TOTAL_SLIDES)


# ============ SLIDE 6: Foundry Agent Factory (6-Step Assembly Line) ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
slide_title_bar(slide, "The Foundry Agent Factory", "6-Step Assembly Line to Production-Ready Agents")

steps = [
    ("1", "Models", "Select LLM\n(GPT-4o, Llama,\nGPT-3.5)", DARK_BLUE),
    ("2", "Customize", "Fine-tune,\ndistillation,\nprompts", GREEN),
    ("3", "Knowledge\n& Tools", "Bing, SharePoint,\nAI Search, Logic\nApps, Functions", PURPLE),
    ("4", "Orchestrate", "Workflows,\ntool calls,\nstate mgmt", ORANGE),
    ("5", "Observe", "Logs, traces,\nevaluations,\nApp Insights", DARK_BLUE),
    ("6", "Trust", "Entra ID, RBAC,\ncontent filters,\nencryption", RED_ACCENT),
]

for i, (num, title, desc, clr) in enumerate(steps):
    x = Inches(0.5 + i * 2.1)
    y = Inches(2.0)

    # Step box
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, y, Inches(1.9), Inches(3.5), CHARCOAL, clr)

    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # Number
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = clr
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)

    # Title
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(8)

    # Description
    p3 = tf.add_paragraph()
    p3.text = desc
    p3.font.size = Pt(10)
    p3.font.color.rgb = LIGHT_GRAY
    p3.alignment = PP_ALIGN.CENTER

    # Arrow between steps
    if i < 5:
        add_arrow_shape(slide, x + Inches(1.95), Inches(3.5), Inches(0.2), Inches(0.3), LIGHT_BLUE)

# Bottom output
output = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                   Inches(2.5), Inches(5.9), Inches(8.3), Inches(0.8), DARK_BLUE, LIGHT_BLUE)
add_text_to_shape(output,
                  "Production-Ready Agent: Reliable, Extensible, Safe to Deploy",
                  font_size=15, bold=True, color=WHITE)

# Arrow from steps to output
add_shape(slide, MSO_SHAPE.DOWN_ARROW,
          Inches(6.3), Inches(5.55), Inches(0.7), Inches(0.35), LIGHT_BLUE)

slide_footer(slide, 6, TOTAL_SLIDES)


# ============ SLIDE 7: Agent 365 Overview ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
slide_title_bar(slide, "Microsoft Agent 365", "Enterprise Agent Lifecycle Management")

# Definition
defn = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                 Inches(0.6), Inches(1.6), Inches(12), Inches(1.0), CHARCOAL, PURPLE)
tf = defn.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("Agent 365 gives each AI agent its own Microsoft Entra Agent ID for identity, lifecycle, "
          "and access management. Agents become observable and manageable in the M365 admin center.")
p.font.size = Pt(13)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

# 5 Pillars
a365_pillars = [
    ("Registry", "Complete view of all\nagents - with Agent ID,\nregistered & shadow", DARK_BLUE),
    ("Access\nControl", "Limit agent access\nto needed resources.\nConditional access", GREEN),
    ("Visualization", "Explore connections\nbetween agents, people\n& data in real time", PURPLE),
    ("Interop", "Equip agents with\napps & data. Connect\nto Work IQ", ORANGE),
    ("Security", "Protect agents from\nthreats. Detect &\nremediate attacks", RED_ACCENT),
]

for i, (title, desc, clr) in enumerate(a365_pillars):
    x = Inches(0.6 + i * 2.5)
    y = Inches(3.0)

    # Pillar box
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, y, Inches(2.2), Inches(2.8), CHARCOAL, clr)

    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = clr
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER

# Stakeholder row
stakeholders = ["IT Admin", "Security", "Developer", "Business", "Info Worker"]
for i, label in enumerate(stakeholders):
    x = Inches(0.6 + i * 2.5)
    s = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                  x, y + Inches(3.1), Inches(2.2), Inches(0.45), clr)
    s.fill.solid()
    s.fill.fore_color.rgb = MID_GRAY
    s.line.fill.background()
    add_text_to_shape(s, label, font_size=10, bold=True, color=WHITE)

slide_footer(slide, 7, TOTAL_SLIDES)


# ============ SLIDE 8: AgentOps Lifecycle ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
slide_title_bar(slide, "The AgentOps Lifecycle", "End-to-End Agent Operations Pipeline")

# Circular lifecycle
phases = [
    ("Design", "Define agent goals,\ninstructions, tools\n& boundaries", DARK_BLUE),
    ("Build", "Develop with SDKs,\ntest locally,\niterate prompts", GREEN),
    ("Deploy", "Push to production\nvia CI/CD, apply\nsecurity policies", PURPLE),
    ("Monitor", "Traces, logs, metrics\nvia App Insights &\nAgent 365 dashboard", ORANGE),
    ("Evaluate", "Run evaluations,\ncollect feedback,\nmeasure quality", RED_ACCENT),
    ("Optimize", "Refine instructions,\nretrain models,\nupdate tools", YELLOW),
]

# Draw as a hexagonal flow
positions = [
    (Inches(5.3), Inches(1.6)),   # Top center
    (Inches(9.0), Inches(2.5)),   # Top right
    (Inches(9.0), Inches(4.5)),   # Bottom right
    (Inches(5.3), Inches(5.4)),   # Bottom center
    (Inches(1.6), Inches(4.5)),   # Bottom left
    (Inches(1.6), Inches(2.5)),   # Top left
]

for i, ((title, desc, clr), (x, y)) in enumerate(zip(phases, positions)):
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, y, Inches(2.6), Inches(1.3), CHARCOAL, clr)

    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = f"{i+1}. {title}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = clr
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER

# Center label
center_label = add_shape(slide, MSO_SHAPE.OVAL,
                         Inches(5.5), Inches(3.3), Inches(2.3), Inches(1.5), DARK_BLUE, LIGHT_BLUE)
add_text_to_shape(center_label, "Continuous\nLoop", font_size=16, bold=True, color=WHITE)

# Arrows between phases (simplified curved arrows using chevrons)
arrow_positions = [
    (Inches(7.9), Inches(2.2)),
    (Inches(10.5), Inches(3.8)),
    (Inches(7.9), Inches(5.5)),
    (Inches(4.2), Inches(5.5)),
    (Inches(1.4), Inches(3.8)),
    (Inches(4.2), Inches(2.2)),
]

for x, y in arrow_positions:
    add_shape(slide, MSO_SHAPE.OVAL, x, y, Inches(0.25), Inches(0.25), LIGHT_BLUE)

slide_footer(slide, 8, TOTAL_SLIDES)


# ============ SLIDE 9: Foundry Agent Architecture ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
slide_title_bar(slide, "Agent Architecture in Foundry", "Three Core Components of Every Agent")

# Three main component blocks
comp_data = [
    ("Model (LLM)", "Powers reasoning &\nlanguage understanding",
     ["GPT-4o", "GPT-4", "GPT-3.5", "Llama", "Custom"], DARK_BLUE),
    ("Instructions", "Define goals, behavior\n& constraints",
     ["Declarative Prompt", "Workflow (YAML)", "Hosted (Code)"], GREEN),
    ("Tools", "Retrieve knowledge\nor take action",
     ["Bing Search", "SharePoint", "AI Search", "Logic Apps", "Functions", "OpenAPI"], PURPLE),
]

for i, (title, desc, items, clr) in enumerate(comp_data):
    x = Inches(0.6 + i * 4.2)
    y = Inches(1.8)

    # Header
    header = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                       x, y, Inches(3.8), Inches(1.0), clr)
    tf = header.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = LIGHT_BLUE
    p2.alignment = PP_ALIGN.CENTER

    # Items list
    for j, item in enumerate(items):
        item_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                             x + Inches(0.2), y + Inches(1.2 + j * 0.5),
                             Inches(3.4), Inches(0.4), CHARCOAL, clr)
        add_text_to_shape(item_box, item, font_size=11, color=LIGHT_GRAY)

# Data flow
add_textbox(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.5),
            "Input (User Prompt / Alert) --> Agent [Model + Instructions + Tools] --> Output (Response / Action)",
            font_size=14, bold=True, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

# Flow diagram at bottom
flow_items = ["User Input", "Agent Processing", "Tool Calls", "Response Output"]
flow_colors = [ORANGE, DARK_BLUE, GREEN, PURPLE]
for i, (label, clr) in enumerate(zip(flow_items, flow_colors)):
    x = Inches(1.0 + i * 3.0)
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, Inches(6.2), Inches(2.5), Inches(0.6), clr)
    add_text_to_shape(box, label, font_size=12, bold=True, color=WHITE)
    if i < 3:
        add_arrow_shape(slide, x + Inches(2.55), Inches(6.35), Inches(0.4), Inches(0.25), LIGHT_BLUE)

slide_footer(slide, 9, TOTAL_SLIDES)


# ============ SLIDE 10: Workflow Orchestration Patterns ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
slide_title_bar(slide, "Workflow Orchestration Patterns", "Multi-Agent Coordination in Foundry")

patterns = [
    ("Sequential", "Pass results from one\nagent to the next in\na defined order",
     "Step-by-step pipelines,\nmulti-stage processing", DARK_BLUE,
     ["Agent A", "Agent B", "Agent C"]),
    ("Group Chat", "Dynamically pass control\nbased on context\nor rules",
     "Escalation, fallback,\nexpert handoff", GREEN,
     ["Agent 1", "Agent 2", "Agent 3"]),
    ("Human-in-Loop", "Ask user questions\nand await input\nto proceed",
     "Approval requests,\nobtaining info", PURPLE,
     ["Agent", "Human", "Agent"]),
]

for i, (title, desc, use_case, clr, agents) in enumerate(patterns):
    y_base = Inches(1.7 + i * 1.9)

    # Pattern label
    label = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                      Inches(0.6), y_base, Inches(2.5), Inches(1.5), clr)
    tf = label.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = LIGHT_BLUE
    p2.alignment = PP_ALIGN.CENTER

    # Agent flow
    for j, agent_name in enumerate(agents):
        x = Inches(3.8 + j * 2.2)
        agent_clr = ORANGE if agent_name == "Human" else CHARCOAL
        a = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                      x, y_base + Inches(0.3), Inches(1.8), Inches(0.8), agent_clr, clr)
        add_text_to_shape(a, agent_name, font_size=11, bold=True, color=WHITE)
        if j < 2:
            add_arrow_shape(slide, x + Inches(1.85),
                            y_base + Inches(0.5), Inches(0.3), Inches(0.25), LIGHT_BLUE)

    # Use case
    uc = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                   Inches(10.5), y_base + Inches(0.15), Inches(2.5), Inches(1.2), CHARCOAL, clr)
    tf = uc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Use Cases:"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = clr
    p.alignment = PP_ALIGN.LEFT
    p2 = tf.add_paragraph()
    p2.text = use_case
    p2.font.size = Pt(9)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.LEFT

slide_footer(slide, 10, TOTAL_SLIDES)


# ============ SLIDE 11: Observability & Evaluation ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
slide_title_bar(slide, "Observability & Evaluation", "Full Visibility into Agent Behavior")

# Left: Observability Stack
obs_header = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                       Inches(0.6), Inches(1.7), Inches(5.8), Inches(0.7), DARK_BLUE)
add_text_to_shape(obs_header, "Observability Stack", font_size=16, bold=True, color=WHITE)

obs_layers = [
    ("Application Insights", "Usage data, request metrics, error rates", GREEN),
    ("Conversation Traces", "Full message history: user-to-agent & agent-to-agent", PURPLE),
    ("Tool Invocation Logs", "Which tools called, inputs/outputs, latency", ORANGE),
    ("Decision Audit Trail", "Why agent chose specific actions & responses", DARK_BLUE),
    ("Content Safety Logs", "Policy violations, filtered content, XPIA attempts", RED_ACCENT),
]

for i, (title, desc, clr) in enumerate(obs_layers):
    y = Inches(2.6 + i * 0.85)
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.6), y, Inches(5.8), Inches(0.7), CHARCOAL, clr)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{title}:  {desc}"
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.LEFT

# Right: Evaluation Framework
eval_header = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(6.9), Inches(1.7), Inches(5.8), Inches(0.7), PURPLE)
add_text_to_shape(eval_header, "Evaluation Framework", font_size=16, bold=True, color=WHITE)

eval_items = [
    ("Cloud Evaluations", "Run automated evals against agent responses"),
    ("Feedback Loops", "Collect user ratings & quality signals"),
    ("A/B Testing", "Compare agent versions side-by-side"),
    ("Drift Detection", "Monitor for quality degradation over time"),
    ("Metrics Dashboard", "Track KPIs: accuracy, latency, satisfaction"),
]

for i, (title, desc) in enumerate(eval_items):
    y = Inches(2.6 + i * 0.85)
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(6.9), y, Inches(5.8), Inches(0.7), CHARCOAL, GREEN)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{title}:  {desc}"
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.LEFT

slide_footer(slide, 11, TOTAL_SLIDES)


# ============ SLIDE 12: Security & Trust ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
slide_title_bar(slide, "Security & Trust", "Enterprise-Grade Agent Protection")

# Security layers diagram
layers = [
    ("Identity & Access", "Microsoft Entra ID, Agent ID,\nRBAC, Conditional Access", DARK_BLUE,
     Inches(0.6), Inches(1.8)),
    ("Content Safety", "Integrated content filters,\nPrompt injection protection (XPIA),\nPolicy-governed outputs", GREEN,
     Inches(4.5), Inches(1.8)),
    ("Network Isolation", "Virtual networks (VNet),\nBring-your-own resources,\nData residency controls", PURPLE,
     Inches(8.4), Inches(1.8)),
    ("Data Protection", "Microsoft Purview integration,\nEncryption at rest & transit,\nAudit logging", ORANGE,
     Inches(0.6), Inches(4.2)),
    ("Threat Protection", "Microsoft Defender integration,\nReal-time threat detection,\nAttack remediation", RED_ACCENT,
     Inches(4.5), Inches(4.2)),
    ("Responsible AI", "RAI validation checks,\nEthical guardrails,\nBias monitoring", YELLOW,
     Inches(8.4), Inches(4.2)),
]

for title, desc, clr, x, y in layers:
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, y, Inches(3.5), Inches(1.8), CHARCOAL, clr)
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = clr
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(6)

    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.LEFT

# Shield icon in center
shield = add_shape(slide, MSO_SHAPE.OVAL,
                   Inches(5.8), Inches(3.65), Inches(1.6), Inches(1.0), DARK_BLUE, LIGHT_BLUE)
add_text_to_shape(shield, "Zero Trust", font_size=13, bold=True, color=WHITE)

slide_footer(slide, 12, TOTAL_SLIDES)


# ============ SLIDE 13: Foundry + Agent 365 Together ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
slide_title_bar(slide, "Foundry + Agent 365: Better Together",
                "Unified AgentOps Across Build & Govern")

# Two sides
# Left: Foundry
foundry_header = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.6), Inches(1.7), Inches(5.5), Inches(0.7), DARK_BLUE)
add_text_to_shape(foundry_header, "Microsoft Foundry (Build & Deploy)", font_size=14, bold=True, color=WHITE)

foundry_items = [
    "Agent Factory: Models + Tools + Instructions",
    "Foundry Agent Service runtime",
    "SDK support: Python, .NET, JS, Java",
    "Workflow orchestration (Sequential, Group Chat, HITL)",
    "Traces, evaluations, App Insights integration",
    "Content safety & prompt injection protection",
]

for i, item in enumerate(foundry_items):
    y = Inches(2.6 + i * 0.55)
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.6), y, Inches(5.5), Inches(0.45), CHARCOAL, DARK_BLUE)
    add_text_to_shape(box, item, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

# Right: Agent 365
a365_header = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(7.2), Inches(1.7), Inches(5.5), Inches(0.7), PURPLE)
add_text_to_shape(a365_header, "Agent 365 (Govern & Operate)", font_size=14, bold=True, color=WHITE)

a365_items = [
    "Entra Agent ID for every agent",
    "Unified registry of all agents (inc. shadow)",
    "Visualization: agents, people, data connections",
    "Access control & conditional access policies",
    "Microsoft Purview data protection",
    "Microsoft Defender threat protection",
]

for i, item in enumerate(a365_items):
    y = Inches(2.6 + i * 0.55)
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(7.2), y, Inches(5.5), Inches(0.45), CHARCOAL, PURPLE)
    add_text_to_shape(box, item, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

# Center bridge
bridge = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                   Inches(4.5), Inches(5.8), Inches(4.3), Inches(1.0), DARK_BLUE, LIGHT_BLUE)
tf = bridge.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "AgentOps = Build (Foundry) + Govern (Agent 365)"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Unified lifecycle from code to compliance"
p2.font.size = Pt(11)
p2.font.color.rgb = LIGHT_BLUE
p2.alignment = PP_ALIGN.CENTER

# Arrows pointing to bridge
add_shape(slide, MSO_SHAPE.DOWN_ARROW,
          Inches(3.0), Inches(5.5), Inches(0.5), Inches(0.35), DARK_BLUE)
add_shape(slide, MSO_SHAPE.DOWN_ARROW,
          Inches(9.8), Inches(5.5), Inches(0.5), Inches(0.35), PURPLE)

slide_footer(slide, 13, TOTAL_SLIDES)


# ============ SLIDE 14: Reference Architecture ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
slide_title_bar(slide, "Enterprise AgentOps Reference Architecture",
                "End-to-End Agent Platform with Microsoft Stack")

# Architecture layers (horizontal bands)
arch_layers = [
    ("User Layer", ["M365 Copilot", "Teams", "Outlook", "Web App", "Custom UI"], DARK_BLUE, Inches(1.7)),
    ("Agent Layer", ["Declarative\nAgents", "Custom Engine\nAgents", "Hosted\nAgents", "Workflow\nAgents"], GREEN, Inches(2.8)),
    ("Platform Layer", ["Foundry Agent Service", "Agent Framework SDK", "Copilot Studio"], PURPLE, Inches(4.0)),
    ("Governance Layer", ["Agent 365 Registry", "Entra Agent ID", "Purview", "Defender"], ORANGE, Inches(5.2)),
    ("Infrastructure", ["Azure AI Services", "Cosmos DB", "AI Search", "Key Vault", "App Insights"], MID_GRAY, Inches(6.3)),
]

for layer_name, items, clr, y in arch_layers:
    # Layer label
    label = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                      Inches(0.3), y, Inches(2.0), Inches(0.8), clr)
    add_text_to_shape(label, layer_name, font_size=11, bold=True, color=WHITE)

    # Items
    item_width = min(Inches(2.0), Inches(10.5 / len(items)))
    for j, item in enumerate(items):
        x = Inches(2.5 + j * (10.5 / len(items)))
        box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                        x, y, item_width, Inches(0.8), CHARCOAL, clr)
        add_text_to_shape(box, item, font_size=9, color=LIGHT_GRAY)

slide_footer(slide, 14, TOTAL_SLIDES)


# ============ SLIDE 15: Getting Started & Next Steps ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
slide_title_bar(slide, "Getting Started with AgentOps", "Your Roadmap to Production Agents")

# Three columns
columns = [
    ("Phase 1: Foundation", [
        "Set up Azure subscription",
        "Create Foundry project",
        "Deploy first model (GPT-4o)",
        "Build a basic agent with SDK",
        "Test locally with playground",
    ], DARK_BLUE),
    ("Phase 2: Production", [
        "Add tools (AI Search, Functions)",
        "Build multi-agent workflows",
        "Enable App Insights tracing",
        "Apply content safety filters",
        "Set up CI/CD pipeline",
    ], GREEN),
    ("Phase 3: Governance", [
        "Enable Agent 365 (Frontier)",
        "Assign Entra Agent IDs",
        "Configure access policies",
        "Enable Purview & Defender",
        "Monitor via M365 admin center",
    ], PURPLE),
]

for i, (title, items, clr) in enumerate(columns):
    x = Inches(0.6 + i * 4.2)
    y = Inches(1.7)

    # Phase header
    header = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                       x, y, Inches(3.8), Inches(0.7), clr)
    add_text_to_shape(header, title, font_size=15, bold=True, color=WHITE)

    # Checklist
    for j, item in enumerate(items):
        iy = y + Inches(0.9 + j * 0.7)
        # Checkbox
        chk = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                        x, iy, Inches(0.35), Inches(0.35), CHARCOAL, clr)
        add_text_to_shape(chk, str(j + 1), font_size=10, bold=True, color=clr)

        add_textbox(slide, x + Inches(0.45), iy, Inches(3.3), Inches(0.5),
                    item, font_size=12, color=WHITE)

    # Arrow between phases
    if i < 2:
        add_arrow_shape(slide, x + Inches(3.9), Inches(2.7), Inches(0.3), Inches(0.35), LIGHT_BLUE)

# Resources footer
resources = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                      Inches(0.6), Inches(6.0), Inches(12), Inches(0.9), CHARCOAL, DARK_BLUE)
tf = resources.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Key Resources"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = LIGHT_BLUE
p.alignment = PP_ALIGN.LEFT
p2 = tf.add_paragraph()
p2.text = ("Foundry Portal: ai.azure.com  |  Agent 365 Docs: learn.microsoft.com/microsoft-agent-365  |  "
           "SDKs: Python, .NET, JS, Java  |  Agent Framework: github.com/microsoft/agents")
p2.font.size = Pt(10)
p2.font.color.rgb = LIGHT_GRAY
p2.alignment = PP_ALIGN.LEFT

slide_footer(slide, 15, TOTAL_SLIDES)


# ========================== SAVE ==========================
output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "docs", "presentations")
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "AgentOps-Microsoft-Foundry-Agent365.pptx")
prs.save(output_path)
print(f"[PASS] Presentation saved to: {output_path}")
print(f"[INFO] {TOTAL_SLIDES} slides generated with diagrams and visuals")
